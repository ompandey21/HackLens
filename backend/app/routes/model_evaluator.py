from fastapi import APIRouter, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from pydantic import BaseModel
import pandas as pd
import numpy as np
from pptx import Presentation
import re
import joblib
import textstat
from sentence_transformers import SentenceTransformer
import onnxruntime as ort
from sklearn.metrics import accuracy_score, precision_score, recall_score, f1_score
import tempfile, os
from app.services.code_executor import execute_code

router = APIRouter()
class CodeRequest(BaseModel):
    language: str
    code: str

TEST_DATA_PATH = r"D:\Projects\Hacklens\backend\app\routes\test.csv"

# ML Model for PPT evaluation
MODEL_PATH = r"D:\Projects\Hacklens\backend\app\ml_model\pitch_model.pkl"
model = joblib.load(MODEL_PATH)
sbert = SentenceTransformer("all-MiniLM-L6-v2")

# Code execution endpoint for JAVA/ CPP/ Python
@router.post("/execute/")
def execute(req: CodeRequest):
    result = execute_code(req.language, req.code)
    return result

# Model Evaluator route
@router.post("/evaluate-model")
async def evaluate_model(file: UploadFile = File(...)):
    """External route version (API endpoint)."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".onnx") as tmp:
            contents = await file.read()
            tmp.write(contents)
            tmp_path = tmp.name

        result = await run_model_evaluation(tmp_path)
        os.remove(tmp_path)
        return JSONResponse(content=result)
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Helper for model evaluator
async def run_model_evaluation(file_path: str) -> dict:
    """Reusable internal function to evaluate ONNX models from file path."""
    if not os.path.exists(TEST_DATA_PATH):
        raise FileNotFoundError("Missing test data CSV at TEST_DATA_PATH.")

    df = pd.read_csv(TEST_DATA_PATH)
    X_test = df.drop(columns=["target"])
    y_test = df["target"].values

    sess = ort.InferenceSession(file_path, providers=["CPUExecutionProvider"])
    input_name = sess.get_inputs()[0].name
    output_name = sess.get_outputs()[0].name

    X_input = X_test.astype(np.float32)
    y_pred = sess.run([output_name], {input_name: X_input.values})[0]
    y_pred = np.argmax(y_pred, axis=1) if y_pred.ndim > 1 else y_pred

    accuracy = accuracy_score(y_test, y_pred)
    precision = precision_score(y_test, y_pred, average='weighted', zero_division=0)
    recall = recall_score(y_test, y_pred, average='weighted', zero_division=0)
    f1 = f1_score(y_test, y_pred, average='weighted', zero_division=0)

    return {
        "accuracy": accuracy,
        "precision": precision,
        "recall": recall,
        "f1_score": f1,
        "final_combined_score": (
            0.4 * accuracy + 0.2 * precision + 0.2 * recall + 0.2 * f1
        ),
    }


# PPT Pitch Deck Evaluation Endpoint
@router.post("/predict-ppt")
async def predict_pitch(file: UploadFile = File(...)):
    try:
        # Save uploaded PPT temporarily
        temp_path = f"temp_{file.filename}"
        with open(temp_path, "wb") as f:
            f.write(await file.read())

        # Extract features
        X, text = extract_features_from_ppt(temp_path)

        # Predict score
        score_10 = float(model.predict(X)[0])
        score_100 = round(score_10 * 10, 2)

        # Generate simple feedback (rule-based)
        feedback = []
        if "problem" not in text.lower():
            feedback.append("Add a clear 'Problem Statement' section.")
        if "solution" not in text.lower():
            feedback.append("Include a concise 'Solution' slide.")
        if "tech" not in text.lower():
            feedback.append("Describe your technology stack clearly.")
        if "future" not in text.lower():
            feedback.append("Mention the 'Future Scope' or roadmap.")
        if "demo" not in text.lower():
            feedback.append("Add a demo/prototype explanation slide.")

        os.remove(temp_path)

        return JSONResponse({
            "score_out_of_10": round(score_10, 2),
            "score_out_of_100": score_100,
            "feedback": feedback if feedback else ["Looks well-balanced! Great job."]
        })

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


# PPT Evaluation Helpers
def extract_features_from_ppt(file_path):
    prs = Presentation(file_path)
    slide_count = len(prs.slides)

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
    joined_text = " ".join(all_text)

    word_count = len(joined_text.split())
    avg_words_per_slide = word_count / slide_count if slide_count else 0
    readability = textstat.flesch_reading_ease(joined_text)

    keywords = {
        "has_problem": int(bool(re.search(r"problem", joined_text, re.I))),
        "has_solution": int(bool(re.search(r"solution", joined_text, re.I))),
        "has_tech": int(bool(re.search(r"tech|technology|stack", joined_text, re.I))),
        "has_future": int(bool(re.search(r"future|scope|next", joined_text, re.I))),
        "has_demo": int(bool(re.search(r"demo|prototype|working", joined_text, re.I))),
    }

    embedding = sbert.encode(joined_text)
    features = [
        slide_count, word_count, avg_words_per_slide, readability,
        keywords["has_problem"], keywords["has_solution"], keywords["has_tech"],
        keywords["has_future"], keywords["has_demo"]
    ] + embedding.tolist()

    return np.array(features).reshape(1, -1), joined_text
