from fastapi import FastAPI, UploadFile, File
from fastapi.responses import JSONResponse
import joblib
import numpy as np
import os
from pptx import Presentation
import re
import textstat
from sentence_transformers import SentenceTransformer

# --- Initialize app ---
app = FastAPI(title="HackLens PPT Evaluator", version="1.0")

@app.get("/")
def root():
    return {
        "message": "✅ HackLens PPT Evaluator API is running! Visit /docs to test file upload."
    }

# --- Load ML model ---
MODEL_PATH = r"C:\Users\Rachit Upadhyay\OneDrive\Desktop\Mini_Project\models\pitch_model.pkl"
model = joblib.load(MODEL_PATH)

# --- Load SBERT model (for embeddings) ---
sbert = SentenceTransformer("all-MiniLM-L6-v2")

# --- Feature extraction (same logic used in notebook) ---
def extract_features_from_ppt(file_path):
    prs = Presentation(file_path)
    slide_count = len(prs.slides)

    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
    joined_text = " ".join(all_text)

    word_count = len(joined_text.split())
    avg_words_per_slide = word_count / slide_count if slide_count else 0
    readability = textstat.flesch_reading_ease(joined_text)

    keywords = {
        "has_problem": int(bool(re.search(r"problem", joined_text, re.I))),
        "has_solution": int(bool(re.search(r"solution", joined_text, re.I))),
        "has_tech": int(bool(re.search(r"tech|technology|stack", joined_text, re.I))),
        "has_future": int(bool(re.search(r"future|scope|next", joined_text, re.I))),
        "has_demo": int(bool(re.search(r"demo|prototype|working", joined_text, re.I))),
    }

    embedding = sbert.encode(joined_text)
    features = [
        slide_count, word_count, avg_words_per_slide, readability,
        keywords["has_problem"], keywords["has_solution"], keywords["has_tech"],
        keywords["has_future"], keywords["has_demo"]
    ] + embedding.tolist()

    return np.array(features).reshape(1, -1), joined_text


# --- Prediction Endpoint ---
@app.post("/predict-ppt")
async def predict_pitch(file: UploadFile = File(...)):
    try:
        # Save uploaded PPT temporarily
        temp_path = f"temp_{file.filename}"
        with open(temp_path, "wb") as f:
            f.write(await file.read())

        # Extract features
        X, text = extract_features_from_ppt(temp_path)

        # Predict score
        score_10 = float(model.predict(X)[0])
        score_100 = round(score_10 * 10, 2)

        # Generate simple feedback (rule-based)
        feedback = []
        if "problem" not in text.lower():
            feedback.append("Add a clear 'Problem Statement' section.")
        if "solution" not in text.lower():
            feedback.append("Include a concise 'Solution' slide.")
        if "tech" not in text.lower():
            feedback.append("Describe your technology stack clearly.")
        if "future" not in text.lower():
            feedback.append("Mention the 'Future Scope' or roadmap.")
        if "demo" not in text.lower():
            feedback.append("Add a demo/prototype explanation slide.")

        os.remove(temp_path)

        return JSONResponse({
            "score_out_of_10": round(score_10, 2),
            "score_out_of_100": score_100,
            "feedback": feedback if feedback else ["Looks well-balanced! Great job."]
        })

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)
